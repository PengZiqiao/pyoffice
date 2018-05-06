from pandas import DataFrame
from pptx import Presentation
from pptx.chart.data import CategoryChartData


class PPT:
    def __init__(self, file):
        self.prs = Presentation(file)

    def __getitem__(self, item):
        """通过页码、shape编号得到指定slide或shape"""
        if isinstance(item, int):
            return self.slides[item]
        elif isinstance(item, str):
            page_idx, shape_idx = item.split()
            return self.slides[int(page_idx)].shapes[int(shape_idx)]
        else:
            raise ValueError('传入key格式有误，请使用int获取指定slide，或"int int"获取指定shape')

    def __setitem__(self, key, value):
        """设置指定shape内容"""
        try:
            shape = self[key]
        except IndexError:
            raise IndexError('无法找到指定shape，请确认page_index或shape_index是否正确')
        else:
            # 填入文字
            if shape.has_text_frame:
                self.replace_text(shape.text_frame, value)

            if isinstance(value, DataFrame):
                # 填写表格
                if shape.has_table:
                    self.fill_table(shape.table, value)
                # 替换图表数据
                elif shape.has_chart:
                    self.replace_chart_data(shape.chart, value)

    @property
    def layouts(self):
        return self.prs.slide_layouts

    @property
    def slides(self):
        return self.prs.slides

    @property
    def _blank_slide_layout(self):
        """获得空白样式的layout"""
        layout_items_count = [len(layout.placeholders)
                              for layout in self.layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return self.layouts[blank_layout_id]

    @staticmethod
    def replace_text(txFrame, text):
        """替换text_frame中的文字"""

        # 只保留第一个paragraph
        for p in txFrame._txBody.p_lst[1:]:
            txFrame._txBody.remove(p)

        # 只保留第一个run
        p = txFrame.paragraphs[0]
        for run in p.runs[1:]:
            p._p.remove(run._r)

        # 替换text
        p.runs[0].text = str(text)

    @staticmethod
    def df2chart_data(df):
        """index转换为chart_data的categories，columns转为chart_data的series"""
        chart_data = CategoryChartData()
        chart_data.categories = df.index.tolist()
        for col in df:
            chart_data.add_series(col, df[col].tolist())
        return chart_data

    def fill_table(self, table, df, index_col=False):
        """用dataframe填写表格"""

        def replace_cell_text(row_idx, col_idx, text):
            cell_text_frame = table.cell(row_idx, col_idx).text_frame
            self.replace_text(cell_text_frame, text)

        # 确定行、列数
        if index_col:
            df.reset_index(inplace=True)
        rows, cols = df.shape

        # 第一行为表头
        columns = df.columns.tolist()
        for col, value in enumerate(columns):
            replace_cell_text(0, col, value)

        # 填入数据
        matrix = df.as_matrix()
        for row in range(rows):
            for col in range(cols):
                replace_cell_text(row + 1, col, matrix[row, col])

    def replace_chart_data(self, chart, df):
        chart_data = self.df2chart_data(df)
        chart.replace_data(chart_data)

    def analyze_layouts(self, output_file='layouts_analyze.pptx'):
        # 遍历每个版式与占位符
        for i, layout in enumerate(self.prs.slide_layouts):
            slide = self.prs.slides.add_slide(layout)

            # 将占位符(placeholders)命名为x-x
            for each in slide.placeholders:
                each.text = f'{i}-{each.placeholder_format.idx}'

            # 是否有标题占位符
            try:
                title = slide.shapes.title
                title.text = f'{i}-标题'
            except AttributeError:
                print(f'[*] layout {i} has no title')

        # 保存
        self.save(output_file)

    def analyze_slides(self, output_file='slides_analyze.pptx'):
        # 遍历每页
        for p, slide in enumerate(self.slides):

            # 将该页原有shape存入字典
            shapes = {}
            for i, shape in enumerate(slide.shapes):
                shapes[i] = shape

            # 为每个shape创建一个对应的label
            for i, shape in shapes.items():
                label = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                label.text = f'{p}-{i}'

        # 保存
        self.save(output_file)

    def save(self, output_file='ouput.pptx'):
        self.prs.save(output_file)
        print(f'[*] {output_file} saved.')
